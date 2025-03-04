from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
from .models import Event, Participant
import pandas as pd
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os
import csv

@staff_member_required
def index(request):
	return render(request, 'index.html')

@staff_member_required
def create(request):
    if request.method == "POST":
        csv_file = request.FILES.get('csv')
        temp = request.FILES.get('template')
        event_name = request.POST.get('event_name')
        
        try:
            dialect = csv.Sniffer().sniff(csv_file.read().decode('latin-1'))
            csv_file.seek(0)  
        except csv.Error:
            messages.error(request, "Invalid CSV file.")
            return redirect('admin:create_certificate')
        
        try:
            presentation = Presentation(temp)
            
        except Exception as e:
            messages.error(request, "Template file is not a valid PowerPoint presentation.")
            return redirect('admin:create_certificate')

        event = Event(
            user=request.user,
            event_name=event_name,
            csv_file=csv_file,
            template=temp
        )
        event.save()

        return redirect(f"/certificate/{event.id}/{event.slug}")

    return render(request, 'admin/certificate/create_event.html')

@staff_member_required
def delete_event(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	if event.user == request.user:
		event.delete()
	return redirect('view_certificate_status')

@staff_member_required
def track(request, id, slug):
	event = Event.objects.filter(slug=slug, id=id).first()
	if event.message:

		return render(request, 'admin/certificate/track.html', {
			'event_name': event.event_name,
			'event_date': event.date,
			'participat_details': Participant.objects.filter(event=event)
			})

	prs = Presentation(event.template)
	st=""
	for slide in prs.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				st = st + shape.text
				st = st + " "
	li = st.split()
	tags = []
	for i in li:
		if i[0] == "<" and i[-1] == ">":
			tags.append(i)

	if request.method == "POST":
		email_col = request.POST.get('emails')
		subject = request.POST.get('subject')
		mess = request.POST.get('mess')
		values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]	
		
		event.email_column = email_col
		event.message = mess
		event.subject = subject
		event.save()

		df=pd.read_csv(event.csv_file)
		df_len=df.shape
		i=0
		
		while i < df_len[0]:
			prs = Presentation(event.template)
			j=""
			if i<9:
				j="00"
			elif i>=9 and i < 99 :
				j="0"
			
			for tag, v_type, value in values:
				for slide in prs.slides:
					for shape in slide.shapes:
						if shape.has_text_frame:
							if(shape.text.find(tag))!=-1:
									text_frame = shape.text_frame
									for paragraph in text_frame.paragraphs:
										for run in paragraph.runs:
											cur_text = run.text
											if v_type == 'text':
												new_text = cur_text.replace(tag, value)
											elif v_type == 'date':
												new_text = cur_text.replace(tag, '/'.join(value.split('-')[::-1]))
											elif v_type == 'csv':
												new_text = cur_text.replace(tag, str(df.loc[i, value]))
											elif v_type == "auto":
												new_text = cur_text.replace(tag, value+"/"+j+str(i+1))
											else:
												pass
											run.text = new_text
			
   
			s_name=df.loc[i,event.email_column].split('@')[0]
			pptx_file_path=s_name+".pptx"
			prs.save(pptx_file_path)
			
			pdf_file_path = pptx_file_path.replace('.pptx', '.pdf')
			presentation = Presentation(pptx_file_path)
			presentation.save(pdf_file_path)
			try:
				mail = EmailMessage(subject,
					f"Hello, {s_name} \n{mess}",
					settings.EMAIL_HOST_USER,
					[df.loc[i,event.email_column]])
				with open(pdf_file_path, 'rb') as f:
					mail.attach('certificate.pdf', f.read(), 'application/pdf')
				mail.send()
				Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
				os.remove(s_name+'.pdf')
				os.remove(s_name+".pptx")
			except:
				Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
				os.remove(s_name+'.pdf')
				os.remove(s_name+".pptx")
			i=i+1

		messages.success(request, "Certificates Sent Successfuly !!")
		return redirect(f"/certificate/{event.id}/{event.slug}")


	return render(request, 'admin/certificate/map_tags_of_template.html',{
		'columns': list(pd.read_csv(event.csv_file).columns),
		'tags': tags,
		})


@staff_member_required
def view_certificate_status(request):
	return render(request, 'admin/certificate/view_certificate_status.html',{
		'events': Event.objects.filter(user=request.user)
		})
